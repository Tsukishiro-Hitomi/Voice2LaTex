"""课件抽取：PDF / PPTX → 每页的文字 + 抽出来的图。

给整合阶段当参考资料用：校正术语和公式、提供章节结构、把图插进笔记。
图片存到 session 的 figures/ 下，LaTeX 里按相对路径引用。
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path

MIN_SIDE = 180          # 太小的当图标/装饰扔掉
MIN_AREA = 90_000       # 约 300x300
MAX_ASPECT = 8.0        # 极端长条一般是分割线/页眉
DRAW_MIN = 12           # 这一页矢量绘图元素多于这个数，才值得整页渲染
PAGE_DPI = 130


@dataclass
class Figure:
    rel: str            # 在 tex 里引用的路径，如 figures/p07_1.png
    page: int
    note: str           # 给模型看的描述

@dataclass
class Slide:
    index: int          # 从 1 开始
    title: str
    text: str
    figures: list[Figure] = field(default_factory=list)


@dataclass
class Deck:
    source: Path
    slides: list[Slide]

    @property
    def figures(self) -> list[Figure]:
        return [f for s in self.slides for f in s.figures]

    def outline(self, per_slide: int = 320, total: int = 9000) -> str:
        """课件正文摘要，塞进 prompt 用。超预算就截断并说明。"""
        parts, used, dropped = [], 0, 0
        for s in self.slides:
            body = " ".join(s.text.split())[:per_slide]
            block = f"[第{s.index}页] {s.title}\n{body}" if s.title else f"[第{s.index}页] {body}"
            if used + len(block) > total:
                dropped += 1
                continue
            parts.append(block)
            used += len(block)
        if dropped:
            parts.append(f"（还有 {dropped} 页因长度限制未列出）")
        return "\n\n".join(parts)

    def figure_manifest(self) -> str:
        if not self.figures:
            return ""
        return "\n".join(f"{f.rel} —— {f.note}" for f in self.figures)


# ---------------- PDF ----------------

def _pdf(src: Path, out: Path) -> list[Slide]:
    import pymupdf

    doc = pymupdf.open(src)
    n_pages = doc.page_count

    # 先统计每个图对象出现在多少页——模板 logo 会到处出现，要排除
    seen: dict[int, int] = {}
    for page in doc:
        for info in page.get_images(full=True):
            seen[info[0]] = seen.get(info[0], 0) + 1
    ubiquitous = {x for x, c in seen.items() if c > max(2, n_pages * 0.3)}

    slides: list[Slide] = []
    saved: set[int] = set()
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        title = lines[0][:60] if lines else ""
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        figs: list[Figure] = []

        for k, info in enumerate(page.get_images(full=True), 1):
            xref = info[0]
            if xref in ubiquitous or xref in saved:
                continue
            try:
                pix = pymupdf.Pixmap(doc, xref)
            except Exception:
                continue
            if pix.n - pix.alpha >= 4:                     # CMYK → RGB
                pix = pymupdf.Pixmap(pymupdf.csRGB, pix)
            w, h = pix.width, pix.height
            if (min(w, h) < MIN_SIDE or w * h < MIN_AREA
                    or max(w, h) / max(1, min(w, h)) > MAX_ASPECT):
                continue
            rel = f"figures/p{i:03d}_{k}.png"
            pix.save(str(out / Path(rel).name))
            saved.add(xref)
            figs.append(Figure(rel, i, f"第{i}页内嵌图{'（' + title + '）' if title else ''}，{w}x{h}"))

        # 没抽到内嵌图但这页画了不少矢量图形 → 整页渲染，图多半是矢量的
        if not figs and len(page.get_drawings()) >= DRAW_MIN:
            rel = f"figures/p{i:03d}_page.png"
            page.get_pixmap(dpi=PAGE_DPI).save(str(out / Path(rel).name))
            figs.append(Figure(rel, i, f"第{i}页整页截图（图为矢量，无法单独抽出）"
                                      f"{'，标题：' + title if title else ''}"))

        slides.append(Slide(i, title, body or text, figs))
    doc.close()
    return slides


# ---------------- PPTX ----------------

def _pptx(src: Path, out: Path) -> list[Slide]:
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def pictures(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from pictures(sh.shapes)
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield sh

    prs = Presentation(str(src))
    slides: list[Slide] = []
    for i, sl in enumerate(prs.slides, 1):
        title = ""
        try:
            if sl.shapes.title is not None:
                title = (sl.shapes.title.text or "").strip()[:60]
        except Exception:
            pass
        texts = [sh.text_frame.text.strip() for sh in sl.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        body = "\n".join(t for t in texts if t != title)

        figs: list[Figure] = []
        for k, sh in enumerate(pictures(sl.shapes), 1):
            blob = sh.image.blob
            try:
                im = Image.open(io.BytesIO(blob))
            except Exception:
                continue
            w, h = im.size
            if (min(w, h) < MIN_SIDE or w * h < MIN_AREA
                    or max(w, h) / max(1, min(w, h)) > MAX_ASPECT):
                continue
            rel = f"figures/s{i:03d}_{k}.png"
            im.convert("RGB").save(out / Path(rel).name)
            figs.append(Figure(rel, i, f"第{i}页图{'（' + title + '）' if title else ''}，{w}x{h}"))
        slides.append(Slide(i, title, body, figs))
    return slides


def extract(src: str | Path, session_dir: Path) -> Deck:
    src = Path(src).expanduser()
    if not src.exists():
        raise FileNotFoundError(f"课件不存在：{src}")
    out = session_dir / "figures"
    out.mkdir(parents=True, exist_ok=True)

    ext = src.suffix.lower()
    if ext == ".pdf":
        slides = _pdf(src, out)
    elif ext == ".pptx":
        slides = _pptx(src, out)
    else:
        raise ValueError(f"课件只支持 .pdf 和 .pptx，收到 {ext}"
                         + ("（.ppt 请先另存为 .pptx 或导出 PDF）" if ext == ".ppt" else ""))
    return Deck(src, slides)
